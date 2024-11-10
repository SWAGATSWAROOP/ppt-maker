from flask import Flask, request, jsonify
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import cloudinary
import cloudinary.uploader
from cloudinary.utils import cloudinary_url
import tempfile
import os
from datetime import datetime
import uuid

app = Flask(__name__)

def create_slide(prs, slide_info):
    """Creates a slide based on the provided information."""
    layout = prs.slide_layouts[1] if slide_info.get("content") else prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)

    # Set title if provided
    if "title" in slide_info:
        title = slide.shapes.title
        title.text = slide_info["title"]

    # Set background color if provided
    if "bg_color" in slide_info:
        rgb_values = slide_info["bg_color"]
        bg_color = RGBColor(rgb_values[0], rgb_values[1], rgb_values[2])
        
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = bg_color

    # Set content text if provided
    if "content" in slide_info:
        content_placeholder = slide.placeholders[1]
        content_placeholder.text = slide_info["content"]

    # Add image placeholders if specified
    if "image_placeholders" in slide_info:
        for placeholder in slide_info["image_placeholders"]:
            left = Inches(placeholder.get("left", 5))
            top = Inches(placeholder.get("top", 2))
            width = Inches(placeholder.get("width", 3))
            height = Inches(placeholder.get("height", 3))
            slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)

def create_presentation(slide_data):
    """Creates a PowerPoint presentation with specified slides."""
    prs = Presentation()
    for slide_info in slide_data:
        create_slide(prs, slide_info)
    return prs

@app.route('/create_ppt', methods=['POST'])
def create_ppt():
    """API route to create a PPT from JSON data and upload to Cloudinary."""
    try:
        cloud_name = request.json.get('cloud_name')
        api_key = request.json.get('api_key')
        api_secret = request.json.get('api_secret')

        if not all([cloud_name, api_key, api_secret]):
            return jsonify({"error": "Cloudinary configuration is incomplete"}), 400
        
        # Configure Cloudinary with dynamic values
        cloudinary.config(
            cloud_name=cloud_name,
            api_key=api_key,
            api_secret=api_secret,
            secure=True
        )

        # Get slide data from request JSON
        slide_data = request.json.get('slides')
        if not slide_data:
            return jsonify({"error": "Slide data is required"}), 400

        # Create presentation from slide data
        prs = create_presentation(slide_data)

        # Generate a unique filename using current time and UUID
        unique_filename = f"{datetime.now().strftime('%Y%m%d%H%M%S')}_{uuid.uuid4()}.pptx"

        # Save presentation to a temporary file
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmpfile:
            prs.save(tmpfile.name)
            tmpfile_path = tmpfile.name

        # Upload to Cloudinary with unique public_id
        upload_result = cloudinary.uploader.upload(tmpfile_path, resource_type="raw", public_id=unique_filename)
        
        # Get downloadable URL
        ppt_url, _ = cloudinary_url(unique_filename, resource_type="raw")

        # Clean up temporary file
        os.remove(tmpfile_path)

        # Return the URL in the response
        return jsonify({"download_url": ppt_url})

    except Exception as e:
        return jsonify({"error": str(e)}), 500

if __name__ == '__main__':
    app.run(debug=True)